"""PPTX to Markdown converter for doc2md."""

import re
import sys
from datetime import datetime
from pathlib import Path

from utils import VERSION
from utils.images import _check_content_sufficient, EMU_PER_PIXEL


def convert_pptx_native(input_path, output_dir, config, args):
    """直接讀取 .pptx，保留投影片結構、標題、表格、圖片、圖表、備忘稿。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        return False

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"Error opening {input_path}: {e}", file=sys.stderr)
        return False

    slide_count = len(prs.slides)
    print(f"Processing: {input_path} ({slide_count} slides)", file=sys.stderr)

    # Image extraction setup
    extract_images_flag = config.get("images", {}).get("extract", True)
    img_subfolder = config.get("images", {}).get("subfolder", "images")
    min_w = config.get("images", {}).get("min_width", 100)
    min_h = config.get("images", {}).get("min_height", 100)

    if extract_images_flag:
        img_dir = output_dir / img_subfolder
        img_dir.mkdir(parents=True, exist_ok=True)

    lines = []
    img_count = 0
    chart_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f'## Slide {slide_num}')
        lines.append('')

        # Extract title
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                lines.append(f'### {title_text}')
                lines.append('')

        # Extract content from shapes
        for shape in slide.shapes:
            # Skip title (already handled)
            if shape is slide.shapes.title:
                continue

            if shape.has_table:
                table = shape.table
                rows_data = []
                for row in table.rows:
                    row_cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                                 for cell in row.cells]
                    rows_data.append(row_cells)

                if rows_data:
                    max_cols = max(len(r) for r in rows_data)
                    for r in rows_data:
                        while len(r) < max_cols:
                            r.append('')

                    header = rows_data[0]
                    table_lines = ['| ' + ' | '.join(header) + ' |']
                    table_lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                    for row in rows_data[1:]:
                        table_lines.append('| ' + ' | '.join(row[:max_cols]) + ' |')

                    lines.append('\n'.join(table_lines))
                    lines.append('')

            elif hasattr(shape, 'has_chart') and shape.has_chart:
                # Extract chart data as markdown table
                try:
                    chart = shape.chart
                    categories = [str(c) for c in chart.categories] if chart.categories else []
                    series_list = []
                    for plot in chart.plots:
                        for series in plot.series:
                            ser_name = series.name if hasattr(series, 'name') and series.name else f'Series {len(series_list)+1}'
                            ser_values = [str(v) if v is not None else '' for v in series.values]
                            series_list.append((ser_name, ser_values))

                    if series_list:
                        if categories:
                            header = ['Category'] + [s[0] for s in series_list]
                            rows = []
                            for i, cat in enumerate(categories):
                                row = [cat]
                                for _, vals in series_list:
                                    row.append(vals[i] if i < len(vals) else '')
                                rows.append(row)
                        else:
                            header = [s[0] for s in series_list]
                            max_len = max(len(s[1]) for s in series_list)
                            rows = []
                            for i in range(max_len):
                                row = []
                                for _, vals in series_list:
                                    row.append(vals[i] if i < len(vals) else '')
                                rows.append(row)

                        max_cols = len(header)
                        table_lines = ['| ' + ' | '.join(header) + ' |']
                        table_lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                        for row in rows:
                            while len(row) < max_cols:
                                row.append('')
                            table_lines.append('| ' + ' | '.join(row[:max_cols]) + ' |')

                        lines.append('\n'.join(table_lines))
                        lines.append('')
                        chart_count += 1
                except Exception:
                    pass

            elif extract_images_flag:
                # Try to detect picture shape
                try:
                    img = shape.image
                    blob = img.blob
                    ext = img.ext
                    if not ext:
                        ext = 'png'

                    # Size filter (shape dimensions in EMU)
                    w_px = shape.width / EMU_PER_PIXEL if shape.width else 0
                    h_px = shape.height / EMU_PER_PIXEL if shape.height else 0
                    if w_px < min_w or h_px < min_h:
                        continue

                    img_count += 1
                    img_filename = f"s{slide_num}_img{img_count}.{ext}"
                    img_path = img_dir / img_filename
                    with open(img_path, 'wb') as f:
                        f.write(blob)

                    lines.append(f"![image](./{img_subfolder}/{img_filename})")
                    lines.append('')
                except (AttributeError, ValueError):
                    # Not a picture shape, try text frame
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                if para.level > 0:
                                    indent = '  ' * (para.level - 1)
                                    lines.append(f'{indent}- {text}')
                                else:
                                    lines.append(text)
                        lines.append('')
                    continue

            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Use indentation level for bullet points
                        if para.level > 0:
                            indent = '  ' * (para.level - 1)
                            lines.append(f'{indent}- {text}')
                        else:
                            lines.append(text)
                lines.append('')

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append(f'> **Notes:** {notes_text}')
                lines.append('')

    if img_count:
        print(f"  Extracted {img_count} image(s)", file=sys.stderr)
    if chart_count:
        print(f"  Extracted {chart_count} chart(s) as table(s)", file=sys.stderr)

    full_md = '\n'.join(lines)
    full_md = re.sub(r'\n{3,}', '\n\n', full_md)
    if not full_md.endswith('\n'):
        full_md += '\n'

    # Frontmatter
    frontmatter = ""
    if not args.no_frontmatter and config.get("frontmatter", {}).get("include", True):
        fm_lines = ["---"]
        fm_lines.append(f"source_file: \"{Path(input_path).name}\"")
        fm_lines.append(f"slides: {slide_count}")
        fm_lines.append(f"converted_date: \"{datetime.now().strftime('%Y-%m-%d %H:%M')}\"")
        fm_lines.append(f"tool_version: \"{VERSION}\"")
        fm_lines.append("---")
        fm_lines.append("")
        frontmatter = '\n'.join(fm_lines)

    final_md = frontmatter + full_md

    # Fallback: if content is too sparse, retry via LibreOffice
    if not _check_content_sufficient(final_md):
        print(f"  Content too sparse ({Path(input_path).name}), falling back to LibreOffice...", file=sys.stderr)
        # Lazy import to avoid circular dependency
        from converters.libreoffice import convert_via_libreoffice
        return convert_via_libreoffice(input_path, output_dir, config, args)

    md_filename = Path(input_path).stem + '.md'
    md_path = output_dir / md_filename
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(final_md)

    print(f"  Output: {md_path}", file=sys.stderr)
    return True
